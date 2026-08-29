#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""《四轴飞控原理分享》v2: 增加机身受力/扭矩分析(开场) + 姿态解算深挖"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import math, os

NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLUE   = RGBColor(0x2E, 0x75, 0xB6)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
RED    = RGBColor(0xC0, 0x00, 0x00)
GRAY   = RGBColor(0x59, 0x59, 0x59)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BORD   = RGBColor(0xB8, 0xC6, 0xDC)
FONT   = "微软雅黑"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def set_run(r, text, size=16, bold=False, color=GRAY, italic=False):
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = FONT
    rPr = r._r.get_or_add_rPr()
    rPr.append(rPr.makeelement(qn('a:ea'), {'typeface': FONT}))

def add_slide(): return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, fill, line=None, round_=False, radius=0.12):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1.2)
    shp.shadow.inherit = False
    return shp

def text_in(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.alignment = align; p.space_after = Pt(4)
        for (t, sz, b, c) in (ln if isinstance(ln, list) else [ln]):
            set_run(p.add_run(), t, sz, b, c)
    return tb

def box_text(shape, lines, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER):
    tf = shape.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(3); tf.margin_top = tf.margin_bottom = Pt(2)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.alignment = align; p.space_after = Pt(2)
        for (t, sz, b, c) in (ln if isinstance(ln, list) else [ln]):
            set_run(p.add_run(), t, sz, b, c)

def box(slide, x, y, w, h, lines, fill=LIGHT, line=BLUE, anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER, round_=True):
    shp = rect(slide, x, y, w, h, fill, line, round_=round_)
    box_text(shp, lines, anchor, align)
    return shp

def _ends(conn, color, w, tail=True, head=False, dash=None):
    conn.line.color.rgb = color; conn.line.width = Pt(w); conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    if dash: ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
    if head: ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if tail: ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))

def arrow(slide, x1, y1, x2, y2, color=NAVY, w=2.0, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    _ends(c, color, w, True, False, dash); return c

def darrow(slide, x1, y1, x2, y2, color=NAVY, w=2.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    _ends(c, color, w, True, True); return c

def gline(slide, x1, y1, x2, y2, color=GRAY, w=1.5, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    _ends(c, color, w, False, False, dash); return c

def circ_arrow(slide, x, y, size, color=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.CIRCULAR_ARROW, Inches(x), Inches(y), Inches(size), Inches(size))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def dot_axis(slide, x, y, label, color=NAVY):
    d = 0.2
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - d/2), Inches(y - d/2), Inches(d), Inches(d))
    c.fill.background(); c.line.color.rgb = color; c.line.width = Pt(1.5); c.shadow.inherit = False
    gline(slide, x - 0.06, y - 0.06, x + 0.06, y + 0.06, color, 1.5)
    gline(slide, x - 0.06, y + 0.06, x + 0.06, y - 0.06, color, 1.5)
    text_in(slide, x + 0.12, y - 0.16, 1.6, 0.3, [(label, 12, True, color)])

PAGE = [0]
def title_bar(slide, title, sub=None):
    PAGE[0] += 1
    rect(slide, 0, 0, 13.333, 0.95, NAVY)
    text_in(slide, 0.55, 0.10, 11.2, 0.8,
            [[(title, 26, True, WHITE)] + ([("   " + sub, 14, False, RGBColor(0xBD,0xD7,0xEE))] if sub else [])],
            anchor=MSO_ANCHOR.MIDDLE)
    text_in(slide, 12.35, 6.98, 0.9, 0.4, [(f"{PAGE[0]:02d}", 12, False, WHITE)], align=PP_ALIGN.RIGHT)
    rect(slide, 0, 0.95, 13.333, 0.045, ORANGE)

def notes(slide, txt): slide.notes_slide.notes_text_frame.text = txt

def table(slide, x, y, w, rows_data, col_w, header_fill=NAVY, size=13, row_h=0.42):
    rows, cols = len(rows_data), len(rows_data[0])
    gt = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(row_h * rows)).table
    for i, cw in enumerate(col_w): gt.columns[i].width = Inches(cw)
    for r in range(rows):
        gt.rows[r].height = Inches(row_h)
        for c in range(cols):
            cell = gt.cell(r, c)
            cell.margin_left = cell.margin_right = Pt(4)
            cell.margin_top = cell.margin_bottom = Pt(1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else (WHITE if r % 2 else LIGHT)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
            set_run(p.add_run(), rows_data[r][c], size, r == 0, WHITE if r == 0 else GRAY)
    return gt

# ============ P1 封面 ============
s = add_slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 5.0, 13.333, 0.05, ORANGE)
text_in(s, 1.0, 1.9, 11.3, 1.4, [("四轴飞控:把原理讲透", 46, True, WHITE)])
text_in(s, 1.0, 3.25, 11.3, 0.8, [("从机身受力到姿态解算:飞机为什么能这么飞,代码为什么这么写", 21, False, RGBColor(0xBD,0xD7,0xEE))])
text_in(s, 1.0, 5.4, 11.3, 1.2, [
    [("一次偏原理的技术分享", 17, True, WHITE)],
    [("主讲:____ · 2026-08 · 配套代码 github.com/xiexiaoan147-cyber/STM32F4_Drone", 14, False, RGBColor(0x8E,0xAA,0xDB))],
])
notes(s, "开场定调:今天从'机身到底受什么力'讲起,一路讲到姿态解算和控制——每一层都回答'为什么'。")

# ============ P2 受力与扭矩① ============
s = add_slide()
title_bar(s, "机身受力与扭矩①", "每个电机给机身两个作用")
# 左: 单电机图解 (俯视)
box(s, 0.55, 1.4, 5.6, 3.5, [], fill=WHITE, line=BORD)
text_in(s, 0.7, 1.48, 4.8, 0.35, [("俯视一个电机:桨盘 + 旋转", 13.5, True, NAVY)])
cx, cy = 3.35, 3.3
pr = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-1.05), Inches(cy-1.05), Inches(2.1), Inches(2.1))
pr.fill.background(); pr.line.color.rgb = GRAY; pr.line.width = Pt(1.5); pr.shadow.inherit = False
circ_arrow(s, cx - 0.45, cy - 0.45, 0.9, BLUE)
text_in(s, cx + 0.42, cy - 0.85, 1.6, 0.4, [("桨:顺时针", 12.5, True, BLUE)])
dot_axis(s, cx, cy, "推力 T:沿机体 z 向上", GREEN)
arrow(s, cx - 1.5, cy + 0.9, cx - 0.35, cy + 0.9, RED, 2.2)
circ_arrow(s, cx - 2.15, cy + 0.28, 0.55, RED)
text_in(s, 0.65, cy + 1.15, 2.2, 0.5, [("反扭矩 M", 12.5, True, RED), ("与桨转向相反", 11, False, RED)])
text_in(s, 0.7, 4.35, 5.3, 0.5, [("推力 T = kt·ω²      反扭矩 M = km·ω²  (都随转速平方变)", 12.5, True, NAVY)])
# 右: 文字说明
text_in(s, 6.5, 1.4, 6.4, 3.5, [
    [("牛顿第三定律的礼物", 16.5, True, NAVY)],
    [("▸ 电机拧着桨转 → 桨拧着空气转", 14, False, GRAY)],
    [("▸ 空气反作用:给桨一个向上的升力,", 14, False, GRAY)],
    [("   同时给机身一个反向扭矩", 14, True, RED)],
    [("", 6, False, GRAY)],
    [("▸ 推力方向永远沿机体 z 轴(垂直机身)", 14, False, GRAY)],
    [("▸ 转速 ×2 → 推力 ×4 (平方律!)", 14, True, BLUE)],
    [("   → 油门与推力不是线性关系,调参时心里有数", 12.5, False, GRAY)],
])
# 下: X 布置
box(s, 0.55, 5.15, 12.35, 1.75, [
    [("X 布置:对角同桨向,相邻反桨向", 14.5, True, NAVY)],
    [("M1 前右(顺) · M2 前左(逆) · M3 后左(顺) · M4 后右(逆)", 13.5, True, GRAY)],
    [("为什么交替?四台全同向 → 反扭矩抵消不掉,机身会一直自旋;交替后正常飞行反扭矩互相抵消,", 13, False, GRAY)],
    [("想要偏航时让'逆桨侧'多转 → 反扭矩出现净差 → 机头转动 (这是 yaw 的唯一来源)", 13, True, BLUE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "一切从单电机开始:它给机身两个作用——沿机身轴的推力+反着桨转向的反扭矩。平方律和'交替桨向'两点是后面混控与yaw控制的物理基础。")

# ============ P3 受力与扭矩② ============
s = add_slide()
title_bar(s, "机身受力与扭矩②", "四个输入怎么变成运动")
box(s, 0.55, 1.3, 5.9, 2.6, [
    [("四个电机 = 四个转速 = 四个控制量", 14.5, True, NAVY)],
    [("总推力  T = T1+T2+T3+T4   (四桨同增减 → 共模)", 13, True, BLUE)],
    [("滚转力矩 τroll  = (左−右)·力臂   (左右差动)", 13, True, BLUE)],
    [("俯仰力矩 τpitch = (前−后)·力臂   (前后差动)", 13, True, BLUE)],
    [("偏航力矩 τyaw   = (逆−顺)·反扭矩差", 13, True, BLUE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 0.55, 4.1, 5.9, 2.7, [
    [("两条动力学方程(全部物理就这两行)", 14, True, NAVY)],
    [("平移: m·a = T·R(q)·ẑ机身 + mg  (牛二)", 14, True, GRAY)],
    [("转动: I·ω̇ = τ  (角动量定理)", 14, True, GRAY)],
    [("", 6, False, GRAY)],
    [("▸ 力矩改变的是'姿态',姿态改变的是", 13, False, GRAY)],
    [("   '推力方向',推力方向才改变'运动'", 13.5, True, RED)],
], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
# 右: 倾斜平移图解 (侧视)
box(s, 6.75, 1.3, 6.15, 3.5, [], fill=WHITE, line=BORD)
text_in(s, 6.9, 1.38, 5.0, 0.35, [("侧视:想往前飞,只能先低头", 13.5, True, NAVY)])
px, py = 8.6, 3.9
# 机身(倾斜线)
ang = math.radians(24)
ex, ey = px + 2.3*math.cos(ang), py - 2.3*math.sin(ang)
rect(s, px - 0.55, py - 0.12, 1.1, 0.24, GRAY, GRAY)
gline(s, px, py, ex, ey, GRAY, 2.2)
m1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ex-0.16), Inches(ey-0.16), Inches(0.32), Inches(0.32))
m1.fill.solid(); m1.fill.fore_color.rgb = ORANGE; m1.line.fill.background(); m1.shadow.inherit = False
# 推力矢量(沿机身轴延长)
arrow(s, ex, ey, ex + 1.5*math.cos(ang), ey - 1.5*math.sin(ang), GREEN, 3.0)
text_in(s, ex + 0.9, ey - 1.35, 1.2, 0.35, [("T 总推力", 13, True, GREEN)])
# 分量
gx2, gy2 = ex + 1.5*math.cos(ang), ey - 1.5*math.sin(ang)
arrow(s, gx2, gy2, gx2, gy2 - 1.1, GREEN, 1.8, dash="dash")
arrow(s, gx2, gy2, gx2 + 1.15, gy2, RED, 2.2)
text_in(s, gx2 - 1.3, gy2 - 1.75, 1.6, 0.4, [("T·cosθ", 12, True, GREEN), ("托住重力", 11, False, GRAY)])
text_in(s, gx2 + 0.15, gy2 + 0.12, 1.7, 0.4, [("T·sinθ → 前飞!", 12.5, True, RED)])
# 重力
arrow(s, px, py + 0.3, px, py + 1.05, GRAY, 2.2)
text_in(s, px + 0.12, py + 0.55, 1.0, 0.3, [("mg", 12.5, True, GRAY)])
text_in(s, 6.9, 4.15, 5.9, 0.55, [("四轴没有水平推进器 → 平移的唯一办法:倾斜,让推力斜着指", 12.5, True, BLUE)])
# 下: 控制级联
box(s, 6.75, 5.05, 6.15, 1.85, [
    [("控制级联(本项目的位置)", 14, True, NAVY)],
    [("", 3, False, GRAY)],
    [("想去哪(位置) → 需要加速度 → 需要姿态角+推力", 12.5, False, GRAY)],
    [("→ 需要三轴力矩 → 电机差动   ←项目实现从这里往下", 12.5, True, BLUE)],
    [("最上层'位置'交给人的手/未来的GPS", 12.5, False, ORANGE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "本页是全场地基:左边四个公式+两条方程就是四轴全部动力学;右图讲欠驱动——想平移只能倾斜,所以'控制姿态=间接控制水平运动'。底部级联图告诉我们项目实现的姿态环在整条链的哪一段。")


# ============ 受力③ X型布局受力图解 ============
s = add_slide()
title_bar(s, "机身受力与扭矩③ X 型布局", "俯视:四个电机怎么分工")

# ---- 左: 悬停基准大图 ----
box(s, 0.55, 1.4, 5.7, 5.5, [], fill=WHITE, line=BORD)
text_in(s, 0.7, 1.5, 5.3, 0.35, [("悬停基准:四桨等速 → 力矩自平衡", 13.5, True, NAVY)])
cx, cy = 3.4, 4.05
arm = 1.12
mots = [("M1", 1, -1, True, "顺桨"), ("M2", -1, -1, False, "逆桨"),
        ("M3", -1, 1, True, "顺桨"), ("M4", 1, 1, False, "逆桨")]
# 机身 X 臂
for sgn in ((-1, -1, 1, 1), (-1, 1, 1, -1)):
    gline(s, cx + sgn[0]*arm*0.82, cy + sgn[1]*arm*0.82,
          cx + sgn[2]*arm*0.82, cy + sgn[3]*arm*0.82, GRAY, 3.0)
for name, dx, dy, cw, lbl in mots:
    mx, my = cx + dx*arm, cy + dy*arm
    # 电机圆
    mc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(mx-0.19), Inches(my-0.19), Inches(0.38), Inches(0.38))
    mc.fill.solid(); mc.fill.fore_color.rgb = RED if cw else BLUE
    mc.line.color.rgb = WHITE; mc.line.width = Pt(1.5); mc.shadow.inherit = False
    box_text(mc, [((name, 11, True, WHITE))])
    # 推力符号(出屏, 朝上)
    dot_axis(s, mx + dx*0.42, my + dy*0.42 - (0.18 if dy<0 else -0.18)*0, "T", GREEN)
    # 桨转向标注
    tx = mx + (0.34 if dx > 0 else -1.5)
    text_in(s, tx, my - 0.5, 1.15, 0.55,
            [[(lbl, 10.5, True, RED if cw else BLUE)], [("反扭" + ("↺" if cw else "↻"), 10, False, RED if cw else BLUE)]]
            if False else [(lbl + (" CW↻" if cw else " CCW↺"), 10.5, True, RED if cw else BLUE)])
# 重心
gc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-0.13), Inches(cy-0.13), Inches(0.26), Inches(0.26))
gc.fill.background(); gc.line.color.rgb = NAVY; gc.line.width = Pt(1.8); gc.shadow.inherit = False
gline(s, cx-0.07, cy-0.07, cx+0.07, cy+0.07, NAVY, 1.6)
gline(s, cx-0.07, cy+0.07, cx+0.07, cy-0.07, NAVY, 1.6)
text_in(s, cx+0.15, cy+0.05, 1.0, 0.3, [("重心", 10.5, True, NAVY)])
text_in(s, 0.7, 5.55, 5.4, 1.25, [
    [("⊙ = 推力(出屏朝上)   ↻↺ = 桨转向", 12, True, GRAY)],
    [("顺桨(M1/M3)反扭矩逆时针,逆桨(M2/M4)反扭矩顺时针", 12, False, GRAY)],
    [("等速时:总推力=mg 且 两两反向的反扭矩抵消 → 不转不歪", 12.5, True, GREEN)],
])

# ---- 右: 三种差动 ----
def mini_x(px, py, title, up, dn, up_names, dn_names, effect, ecol):
    """up/dn: 电机坐标集合 [(dx,dy),..] 加速/减速; effect: 效果文字"""
    box(s, px, py, 6.6, 1.7, [], fill=WHITE, line=BORD)
    text_in(s, px + 0.1, py + 0.06, 3.6, 0.32, [(title, 12.5, True, NAVY)])
    mx0, my0, rr = px + 0.95, py + 0.95, 0.52
    for sgn in ((-1, -1, 1, 1), (-1, 1, 1, -1)):
        gline(s, mx0 + sgn[0]*rr*0.75, my0 + sgn[1]*rr*0.75,
              mx0 + sgn[2]*rr*0.75, my0 + sgn[3]*rr*0.75, GRAY, 2.0)
    for name, dx, dy, cw, lbl in mots:
        mx, my = mx0 + dx*rr, my0 + dy*rr
        if (dx, dy) in up: col, mk = GREEN, "↑"
        elif (dx, dy) in dn: col, mk = RED, "↓"
        else: col, mk = GRAY, "·"
        mc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(mx-0.1), Inches(my-0.1), Inches(0.2), Inches(0.2))
        mc.fill.solid(); mc.fill.fore_color.rgb = col; mc.line.fill.background(); mc.shadow.inherit = False
        text_in(s, mx + dx*0.14 - 0.1, my + dy*0.14 - 0.34, 0.3, 0.28, [(mk, 12, True, col)])
    text_in(s, px + 1.85, py + 0.4, 2.6, 1.15, [
        [("加速↑: " + up_names, 11.5, True, GREEN)],
        [("减速↓: " + dn_names, 11.5, True, RED)],
    ])
    eff_lines = [(ln, 12.5, True, ecol) for ln in effect.split("\n")]
    text_in(s, px + 4.55, py + 0.55, 2.0, 0.9, eff_lines, align=PP_ALIGN.LEFT)

text_in(s, 7.0, 1.42, 5.9, 0.35, [("想动哪根轴,就让对应一侧差动", 13.5, True, NAVY)])
mini_x(7.0, 1.85, "+ Roll(右滚):左侧快右侧慢",
       [(-1, -1), (-1, 1)], [(1, -1), (1, 1)], "M2·M3", "M1·M4", "机身右倾\n→向右平移", BLUE)
mini_x(7.0, 3.65, "+ Pitch(抬头):后侧快前侧慢",
       [(-1, 1), (1, 1)], [(-1, -1), (1, -1)], "M3·M4", "M1·M2", "机头上抬\n→向前飞", BLUE)
mini_x(7.0, 5.45, "+ Yaw(右转):逆桨快顺桨慢",
       [(-1, -1), (1, 1)], [(-1, 1), (1, -1)], "M2·M4", "M1·M3", "反扭矩净差\n→机头右转", BLUE)
notes(s, "把①②页的公式画成图:左边悬停时自平衡;右边三行分别对应三个力矩——注意 roll/pitch 靠'位置差动',yaw 靠'桨向差动'(反扭矩),这就是混控矩阵的物理来源。")


# ============ P4 分享脉络 ============
s = add_slide()
title_bar(s, "分享脉络")
items = [
    ("00", "飞行力学", "受力·扭矩·为什么控制姿态就是控制运动 (刚讲完)"),
    ("01", "IMU 工作原理", "陀螺仪 · 加速度计 · 为什么必须融合 (图解)"),
    ("02", "姿态解算", "四元数 · Mahony 框架·误差本质·数值细节 (重点)"),
    ("03", "控制与实现", "PID · 串级 · 混控 · 状态机 · 安全链"),
]
y = 1.9
for num, t, d in items:
    box(s, 1.6, y, 1.0, 0.85, [(num, 24, True, ORANGE)], fill=NAVY, line=NAVY)
    box(s, 2.8, y, 9.0, 0.85, [[(t + "   ", 19, True, NAVY), (d, 14, False, GRAY)]],
        fill=LIGHT, line=RGBColor(0xD9,0xE2,0xF3), align=PP_ALIGN.LEFT)
    y += 1.12
notes(s, "回顾路线:刚讲完飞机为什么能飞,接下来讲'我们怎么知道它现在的姿态'(IMU+解算),最后讲'怎么控制'(PID)。")

# ============ P5 陀螺仪 ============
s = add_slide()
title_bar(s, "IMU 原理① 陀螺仪", "旋转让振动的东西受到侧向力")

def gyro_panel(px, title, rotated):
    box(s, px, 1.45, 3.0, 3.35, [], fill=WHITE, line=BORD)
    text_in(s, px + 0.1, 1.52, 2.8, 0.35, [(title, 13.5, True, NAVY)], align=PP_ALIGN.CENTER)
    cx, cy = px + 1.5, 3.15
    off = 0.3 if rotated else 0.0
    gline(s, px + 0.35, cy, cx - 0.42 + off, cy, GRAY, 1.5)
    gline(s, cx + 0.42 + off, cy, px + 2.65, cy, GRAY, 1.5)
    for k in (-0.18, 0, 0.18):
        gline(s, px + 0.35, cy + k - 0.09, px + 0.35, cy + k + 0.09, BLUE, 1.2)
        gline(s, px + 2.65, cy + k - 0.09, px + 2.65, cy + k + 0.09, BLUE, 1.2)
    m = rect(s, cx - 0.42 + off, cy - 0.42, 0.84, 0.84, RGBColor(0xFD,0xF0,0xE6), ORANGE)
    box_text(m, [("m", 14, True, ORANGE)])
    text_in(s, px + 0.1, cy + 0.52, 2.8, 0.3, [("质量块", 11, False, GRAY)], align=PP_ALIGN.CENTER)
    darrow(s, cx, cy - 1.28, cx, cy - 0.62, BLUE, 1.8)
    text_in(s, cx + 0.08, cy - 1.18, 1.0, 0.5, [("v", 13, True, BLUE), ("振动", 10.5, False, BLUE)])
    if rotated:
        circ_arrow(s, px + 2.2, 1.62, 0.55, NAVY)
        text_in(s, px + 2.0, 2.2, 0.9, 0.3, [("ω", 14, True, NAVY)])
        arrow(s, cx + 0.5 + off, cy - 0.9, cx + 0.5 + off, cy - 0.15, RED, 2.5)
        text_in(s, cx + 0.62 + off, cy - 0.85, 1.2, 0.6, [("F 科氏力", 12, True, RED), ("侧向推", 10.5, False, RED)])
gyro_panel(0.55, "① 静止:质量块被驱动上下振动", False)
gyro_panel(3.75, "② 旋转:振动+旋转 → 侧向科氏力", True)

box(s, 0.55, 5.05, 6.2, 1.85, [
    [("输出 = 角速度 ω (rad/s)", 15, True, NAVY)],
    [("▸ 电容极板检测质量块位移 → ∝ 角速度", 13, False, GRAY)],
    [("▸ F = −2m (ω × v):v 已知,测 F 反解 ω", 13, True, BLUE)],
], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
text_in(s, 7.15, 1.4, 5.75, 5.5, [
    [("三个必须记住的特点", 17, True, NAVY)],
    [("", 4, False, GRAY)],
    [("① 它测角速度,不测角度!", 15.5, True, RED)],
    [("   角度要自己积分:θ = ∫ω·dt", 14, False, GRAY)],
    [("", 6, False, GRAY)],
    [("② 有零偏 bias:静止也输出'假角速度'", 15.5, True, ORANGE)],
    [("   每颗不同(1~3°/s),必须上电校准", 13.5, False, GRAY)],
    [("", 6, False, GRAY)],
    [("③ 积分误差随时间累积", 15.5, True, ORANGE)],
    [("   1°/s 没校准 → 姿态每分钟漂 60°", 13.5, False, GRAY)],
    [("", 10, False, GRAY)],
    [("结论:陀螺'快但会漂' → 需要绝对基准纠偏", 15, True, NAVY)],
])
notes(s, "图解:静止时质量块只是上下振动;一旋转就被科氏力推向一侧,梳齿电容变化——测偏移即测角速度。三个特点层层递进,引出加速度计。")

# ============ P6 加速度计 ============
s = add_slide()
title_bar(s, "IMU 原理② 加速度计", "先想一个电梯问题")

def elevator(px, title, reading, acc_dir, read_color):
    box(s, px, 1.45, 1.95, 3.1, [], fill=WHITE, line=BORD)
    text_in(s, px + 0.05, 1.52, 1.85, 0.35, [(title, 13.5, True, NAVY)], align=PP_ALIGN.CENTER)
    rect(s, px + 0.42, 1.95, 1.1, 1.85, LIGHT, BLUE)
    hd = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px + 0.83), Inches(2.1), Inches(0.28), Inches(0.28))
    hd.fill.solid(); hd.fill.fore_color.rgb = ORANGE; hd.line.fill.background(); hd.shadow.inherit = False
    rect(s, px + 0.78, 2.44, 0.38, 0.75, RGBColor(0xFD,0xF0,0xE6), ORANGE)
    rect(s, px + 0.55, 3.42, 0.84, 0.16, GRAY, GRAY)
    text_in(s, px + 0.1, 3.7, 1.75, 0.6, [(reading, 13.5, True, read_color)], align=PP_ALIGN.CENTER)
    if acc_dir == 'up':
        arrow(s, px + 1.72, 2.75, px + 1.72, 2.05, RED, 2.2)
        text_in(s, px + 1.6, 2.78, 0.5, 0.3, [("a", 13, True, RED)])
    elif acc_dir == 'down':
        arrow(s, px + 1.72, 2.1, px + 1.72, 2.8, GREEN, 2.2)
        text_in(s, px + 1.55, 2.15, 0.7, 0.3, [("自由落体", 10.5, True, GREEN)])
elevator(0.55, "静止", "秤 = 1g", None, GRAY)
elevator(2.75, "向上加速", "秤 > 1g 超重", 'up', RED)
elevator(4.95, "自由下落", "秤 = 0 失重", 'down', GREEN)

text_in(s, 0.55, 4.75, 6.35, 2.3, [
    [("关键洞察", 16, True, NAVY)],
    [("▸ 秤测的从来不是'加速度',是'支撑力'", 14.5, True, BLUE)],
    [("▸ 加速度计就是那个'秤':它测的是比力", 14.5, True, BLUE)],
    [("      f = a − g   (specific force)", 15, True, NAVY)],
    [("▸ 静止时 a=0 → 输出恰好是 −g(重力!)", 14.5, False, GRAY)],
    [("   → 免费获得'永远指向地心'的参考", 14.5, True, GREEN)],
])
box(s, 7.15, 1.45, 5.75, 2.5, [
    [("于是加速度计 = 倾角计", 15.5, True, NAVY)],
    [("", 4, False, GRAY)],
    [("roll  = atan2( ay , az )", 15.5, True, BLUE)],
    [("pitch = atan2( −ax , √(ay²+az²) )", 15.5, True, BLUE)],
    [("", 4, False, GRAY)],
    [("板子一倾斜,重力在三轴的分量就变", 13, False, GRAY)],
    [("→ 解三角函数就能反推出倾角", 13, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 7.15, 4.2, 5.75, 2.7, [
    [("但它有个致命弱点", 15.5, True, RED)],
    [("▸ 它分不清'重力'和'运动加速度'!", 14, True, RED)],
    [("▸ 电机一转、机身一晃,读数立刻被污染", 13.5, False, GRAY)],
    [("▸ 转弯的离心加速度也会被当成'倾斜'", 13.5, False, GRAY)],
    [("", 6, False, GRAY)],
    [("结论:加计'准但脏' → 静止是天使,机动是噪音", 14.5, True, NAVY)],
], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "电梯类比讲比力:超重失重时秤变了,重力没变——加速度计就是秤。静止时输出即重力方向,所以能当倾角计;一动就脏。两个结论为融合做铺垫。")

# ============ P7 为什么融合 ============
s = add_slide()
title_bar(s, "IMU 原理③ 为什么融合", "一个管快,一个管准")
gline(s, 1.1, 4.55, 1.1, 1.7, GRAY, 1.8)
arrow(s, 1.1, 4.55, 10.3, 4.55, GRAY, 1.8)
text_in(s, 0.25, 1.62, 0.85, 0.3, [("可信度", 12, True, GRAY)])
text_in(s, 1.0, 4.62, 1.2, 0.3, [("低频", 12.5, True, GRAY)])
text_in(s, 9.3, 4.62, 1.2, 0.3, [("高频", 12.5, True, GRAY)])
text_in(s, 5.2, 4.62, 2.4, 0.3, [("信号的频率 →", 12.5, False, GRAY)])
x0, x1 = 1.1, 10.0
def curve_pts(fn, n=28):
    return [(x0 + (x1 - x0) * i / n, 4.55 - fn(i / n) * 2.6) for i in range(n + 1)]
acc_pts = curve_pts(lambda t: math.exp(-3.2 * t) * 0.85 + 0.06)
gyr_pts = curve_pts(lambda t: 0.22 + 0.68 * (1 - math.exp(-9 * t)))
for i in range(len(acc_pts) - 1):
    gline(s, acc_pts[i][0], acc_pts[i][1], acc_pts[i+1][0], acc_pts[i+1][1], GREEN, 2.8)
for i in range(len(gyr_pts) - 1):
    gline(s, gyr_pts[i][0], gyr_pts[i][1], gyr_pts[i+1][0], gyr_pts[i+1][1], BLUE, 2.8)
text_in(s, 1.35, 1.85, 2.6, 0.35, [("加速度计", 14, True, GREEN)])
text_in(s, 7.9, 1.95, 2.0, 0.35, [("陀螺仪", 14, True, BLUE)])
gline(s, 5.05, 1.75, 5.05, 4.55, RED, 1.5, dash="dash")
text_in(s, 4.35, 1.45, 2.3, 0.3, [("交界(截止频率)", 12, True, RED)])
box(s, 1.3, 5.35, 4.2, 1.5, [
    [("互补滤波 = 频域分工", 14.5, True, NAVY)],
    [("姿态 = 低通(加计) + 高通(陀螺)", 14, True, BLUE)],
    [("两条曲线拼起来 = 全带宽可信", 13, False, GRAY)],
], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 5.85, 5.35, 7.05, 1.5, [
    [("落到这个项目上", 14.5, True, NAVY)],
    [("▸ 陀螺积分给'动态',加计长期'拽回'", 13, False, GRAY)],
    [("▸ Mahony 就是这种融合的经典实现(后面详拆)", 13, True, BLUE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "这张图是姿态解算的动机:绿色加计低频可信,蓝色陀螺高频可信。互补滤波=各干各擅长的频段。")

# ============ P8 姿态表示 ============
s = add_slide()
title_bar(s, "姿态怎么表示", "先画两个坐标系")
box(s, 0.55, 1.4, 6.0, 4.1, [], fill=WHITE, line=BORD)
text_in(s, 0.7, 1.48, 3.4, 0.35, [("世界系(固定不动): 东-北-地", 13, True, GRAY)])
wx, wy = 2.15, 4.15
arrow(s, wx, wy, wx + 1.35, wy, GRAY, 2.0); text_in(s, wx + 1.4, wy - 0.16, 0.5, 0.3, [("y东", 12.5, True, GRAY)])
arrow(s, wx, wy, wx, wy - 1.35, GRAY, 2.0); text_in(s, wx + 0.06, wy - 1.62, 0.5, 0.3, [("x北", 12.5, True, GRAY)])
dot_axis(s, wx, wy, "z地(出屏)", GRAY)
bx, by = 4.75, 3.35
text_in(s, 3.7, 1.48, 2.8, 0.35, [("机体系(跟着飞机转)", 13, True, NAVY)])
ang = math.radians(32)
ax1, ay1 = bx + 1.35 * math.cos(ang), by - 1.35 * math.sin(ang)
ax2, ay2 = bx + 0.45 * math.cos(ang), by + 1.15 * math.sin(ang)
arrow(s, bx, by, ax1, ay1, NAVY, 2.0); text_in(s, ax1 + 0.05, ay1 - 0.1, 0.6, 0.3, [("x前", 12.5, True, NAVY)])
arrow(s, bx, by, ax2, ay2, NAVY, 2.0); text_in(s, ax2 + 0.08, ay2 + 0.02, 0.6, 0.3, [("y右", 12.5, True, NAVY)])
dot_axis(s, bx, by, "z下(出屏)", NAVY)
for d1, d2 in ((-0.55, -0.55), (0.55, 0.55)):
    gline(s, bx + d1, by + d1, bx + d2, by + d2, ORANGE, 2.6)
for ddx, ddy in ((-0.62, -0.62), (0.62, -0.62), (-0.62, 0.62), (0.62, 0.62)):
    mc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(bx + ddx - 0.11), Inches(by + ddy - 0.11), Inches(0.22), Inches(0.22))
    mc.fill.solid(); mc.fill.fore_color.rgb = ORANGE; mc.line.fill.background(); mc.shadow.inherit = False
circ_arrow(s, 3.15, 2.6, 0.6, RED)
text_in(s, 3.5, 1.85, 3.2, 0.6, [("姿态 q =", 14, True, RED), ("这两个坐标系的旋转", 13, False, RED)])
text_in(s, 0.7, 5.0, 5.8, 0.45, [("传感器读数在机体系;重力在世界系 z 轴 → 解算=找两者关系", 12.5, True, BLUE)])
table(s, 7.0, 1.4, 5.9, [
    ["表示方法", "参数", "一句话评价"],
    ["欧拉角", "3", "直观;有万向节锁;适合给人看"],
    ["旋转矩阵", "9", "无奇异;但参数冗余运算大"],
    ["四元数", "4", "无奇异;只有乘加;工程首选"],
], [1.35, 0.75, 3.8], size=12.5, row_h=0.52)
box(s, 7.0, 4.0, 5.9, 2.9, [
    [("万向节锁,一句话版", 14.5, True, RED)],
    [("▸ 欧拉角规定按 Z→Y→X 三次转", 13, False, GRAY)],
    [("▸ 抬头到 90° 那一刻,第一次转的", 13, False, GRAY)],
    [("   偏航和第三次转的横滚变成同一根轴", 13, True, RED)],
    [("▸ 三个自由度只剩两个 → 姿态'卡死'", 13, True, RED)],
    [("", 6, False, GRAY)],
    [("四元数/矩阵'一次转到位',没有这个坑", 13.5, True, GREEN)],
], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "姿态=两个坐标系之间的旋转,这句话立住了后面全通。万向节锁就讲'两次旋转撞到同一根轴'。")

# ============ P9 四元数 ============
s = add_slide()
title_bar(s, "四元数怎么用", "别怕,工程上只用三招")
box(s, 0.55, 1.4, 6.0, 2.75, [], fill=WHITE, line=BORD)
text_in(s, 0.7, 1.48, 4.5, 0.35, [("直觉:单位四元数 = 绕某根轴转一个角", 13.5, True, NAVY)])
qx, qy = 2.1, 3.55
arrow(s, qx, qy, qx + 1.7, qy - 1.15, NAVY, 2.5)
text_in(s, qx + 1.35, qy - 1.45, 1.6, 0.35, [("转轴 n", 13, True, NAVY)])
circ_arrow(s, qx + 1.35, qy - 0.75, 0.62, ORANGE)
text_in(s, qx + 0.55, qy - 0.35, 1.2, 0.35, [("转角 θ", 12.5, True, ORANGE)])
box(s, 1.2, 1.95, 4.6, 0.62, [("q = [ cos(θ/2),  n·sin(θ/2) ]", 15, True, BLUE)], fill=LIGHT, line=BLUE)
box(s, 0.55, 4.45, 6.0, 2.45, [
    [("工程上只需记住三件事", 15, True, NAVY)],
    [("① q⊗p = 两个旋转叠加(先q后p)", 13.5, True, GRAY)],
    [("② |q| 必须保持 = 1(每步归一化)", 13.5, True, GRAY)],
    [("③ 共轭 q* = 反着转(坐标系互转用)", 13.5, True, GRAY)],
    [("", 6, False, GRAY)],
    [("内部用四元数算,给人看时转欧拉角", 13.5, True, GREEN)],
], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 6.95, 1.4, 5.95, 2.45, [
    [("唯一的核心方程:姿态更新", 15, True, NAVY)],
    [("q̇ = ½ q ⊗ [0, ω]", 19, True, BLUE)],
    [("每个周期:  q ← q + q̇·dt , 再归一化", 14.5, True, BLUE)],
    [("▸ 陀螺的角速度,就这样'转'进姿态里", 13, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 6.95, 4.15, 5.95, 2.75, [
    [("输出:转成欧拉角给人看", 15, True, NAVY)],
    [("roll  = atan2( 2(q0q1+q2q3), 1−2(q1²+q2²) )", 13, True, BLUE)],
    [("pitch = asin( 2(q0q2−q3q1) )", 13, True, BLUE)],
    [("yaw   = atan2( 2(q0q3+q1q2), 1−2(q2²+q3²) )", 13, True, BLUE)],
    [("", 6, False, GRAY)],
    [("三个易错点:ω 用 rad/s 且减零偏;dt 用实测;", 12.5, False, ORANGE)],
    [("asin 参数钳位防 NaN", 12.5, False, ORANGE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "降低恐惧:四元数工程用法三招。左图轴角直觉,右盒是项目真正在跑的全部公式。")

# ============ P10 Mahony 框图 ============
s = add_slide()
title_bar(s, "Mahony 互补滤波", "总框图:陀螺主内,加计纠偏")
box(s, 0.6, 1.45, 1.75, 0.85, [("陀螺 ω", 14, True, GRAY), ("rad/s·去零偏", 11, False, GRAY)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
box(s, 2.85, 1.55, 0.55, 0.65, [("⊕", 17, True, NAVY)], fill=WHITE, line=NAVY)
box(s, 3.85, 1.45, 2.15, 0.85, [("积分更新", 13.5, True, NAVY), ("q ← q+½q⊗ω'·dt", 11.5, False, GRAY)], fill=LIGHT, line=BLUE)
box(s, 6.55, 1.45, 1.7, 0.85, [("四元数 q", 13.5, True, NAVY), ("每步归一化", 11, False, GRAY)], fill=LIGHT, line=BLUE)
box(s, 8.85, 1.45, 1.7, 0.85, [("欧拉角", 13.5, True, GRAY), ("roll/pitch/yaw", 11, False, GRAY)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
box(s, 6.55, 3.35, 1.7, 0.85, [("重力预测", 13, True, NAVY), ("v = R(q)ᵀ·[0,0,1]", 10.5, False, GRAY)], fill=LIGHT, line=BLUE)
box(s, 4.6, 3.35, 0.55, 0.85, [("×", 16, True, RED)], fill=WHITE, line=RED)
box(s, 0.6, 3.35, 1.75, 0.85, [("加计 a", 14, True, GRAY), ("归一化为单位向量", 10.5, False, GRAY)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
box(s, 2.4, 4.85, 2.5, 0.9, [("PI 修正", 14, True, NAVY), ("ω' = ω + Kp·e + Ki·∫e", 12, True, BLUE)], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE)
arrow(s, 2.35, 1.87, 2.85, 1.87, GREEN)
arrow(s, 3.4, 1.87, 3.85, 1.87)
arrow(s, 6.0, 1.87, 6.55, 1.87)
arrow(s, 8.25, 1.87, 8.85, 1.87, GREEN)
arrow(s, 7.4, 2.3, 7.4, 3.35)
arrow(s, 6.55, 3.77, 5.15, 3.77)
arrow(s, 2.35, 3.77, 4.6, 3.77, GREEN)
arrow(s, 4.87, 4.2, 4.3, 4.85, RED)
arrow(s, 3.65, 5.3, 3.1, 5.3, ORANGE)
arrow(s, 3.1, 5.3, 3.1, 2.2, ORANGE)
text_in(s, 5.25, 2.95, 1.9, 0.35, [("误差 e", 12.5, True, RED)], align=PP_ALIGN.CENTER)
text_in(s, 0.35, 5.35, 2.9, 0.4, [("修正量回注陀螺", 11.5, True, ORANGE)])
box(s, 10.9, 3.35, 2.0, 2.4, [
    [("有效性门限", 12.5, True, NAVY)],
    [("|a| 超出", 11.5, False, GRAY)],
    [("[0.5,20]m/s²", 11.5, False, GRAY)],
    [("跳过修正", 11.5, True, RED)],
    [("(机动时只信陀螺)", 11, False, GRAY)],
], fill=LIGHT, line=BLUE)
notes(s, "全景图:主链路是上面一行;加计从左进来,经叉积得误差,PI后'拧'陀螺——注意不是直接改q,而是修正角速度再积分,这是Mahony优雅之处。比喻:陀螺反应快但越走越歪,加计慢半拍但永远知道下方在哪。")

# ============ P11 Mahony 四步 ============
s = add_slide()
title_bar(s, "Mahony 四步拆解", "每步一行代码级")
steps = [
    ("① 重力预测", "由当前 q 算'重力应在机体系哪个方向'",
     "v = R(q)ᵀ·[0,0,1]",
     "vx = q1q3 − q0q2   vy = q0q1 + q2q3   vz = q0² − ½ + q3²"),
    ("② 误差 = 叉积", "实测重力方向 × 预测方向",
     "e = a_norm × v",
     "方向=转回去的轴, 大小=sin夹角 (两个都是单位向量)"),
    ("③ PI 修正", "误差换算成角速度修正量",
     "ω' = ω + Kp·e + Ki·∫e·dt",
     "Kp:立即拽回(收敛快但抖)  Ki:消恒值零偏"),
    ("④ 积分+归一化", "修正后的角速度进积分器",
     "q ← q + ½·q⊗ω'·dt ; q ← q/|q|",
     "归一化防数值发散"),
]
y = 1.3
for t, d, f, g in steps:
    box(s, 0.55, y, 1.85, 1.05, [(t, 13.5, True, WHITE)], fill=NAVY, line=NAVY)
    box(s, 2.55, y, 3.6, 1.05, [[(d, 11.5, False, GRAY)], [(f, 13.5, True, BLUE)]], fill=LIGHT, line=BLUE)
    box(s, 6.3, y, 6.55, 1.05, [(g, 12, False, GRAY)], fill=WHITE, line=RGBColor(0xD9,0xE2,0xF3))
    y += 1.25
box(s, 0.55, 6.35, 12.3, 0.9, [
    [("调参指引", 13, True, NAVY)],
    [("Kp 从 0.5 起:静止收敛 1~2s 为宜,太快抖/太慢肉;Ki 取 Kp 的 1/10。验证:手转任意角度回平,roll/pitch 应秒级归零", 12.5, False, GRAY)],
], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "四步循环讲,每步配一句人话。第②步是灵魂,下一页专门拆。")

# ============ P12 误差叉积深挖 ============
s = add_slide()
title_bar(s, "深挖① 误差叉积为什么能用", "几何 + 来路 + PI 的物理意义")
# 左: 几何图 (单位球上两向量)
box(s, 0.55, 1.35, 5.7, 3.6, [], fill=WHITE, line=BORD)
text_in(s, 0.7, 1.43, 5.2, 0.35, [("两个单位向量:预测的重力 vs 实测的重力", 13, True, NAVY)])
vx0, vy0 = 3.4, 4.3
axv, ayv = 1.45, 2.35
bxv, byv = 4.15, 2.0
arrow(s, vx0, vy0, axv, ayv, GREEN, 2.6)
text_in(s, axv - 0.75, ayv - 0.34, 1.2, 0.3, [("v̂ 预测", 12.5, True, GREEN)])
arrow(s, vx0, vy0, bxv, byv, ORANGE, 2.6)
text_in(s, bxv + 0.08, byv - 0.32, 1.2, 0.3, [("â 实测", 12.5, True, ORANGE)]
)
# 夹角弧 + 轴
circ_arrow(s, vx0 - 0.95, vy0 - 1.15, 0.75, RED)
text_in(s, 1.9, 4.35, 1.3, 0.4, [("夹角 θe", 12, True, RED)])
dot_axis(s, 3.3, 2.85, "e 的方向", RED)
text_in(s, 0.7, 4.55, 5.4, 0.4, [("e = â × v̂ :|e| = sin θe ,方向 = 把 v̂ 转到 â 的那根轴", 12.5, True, BLUE)])
box(s, 0.55, 5.15, 5.7, 1.75, [
    [("为什么这是'姿态误差'?", 13.5, True, NAVY)],
    [("▸ 两个重力方向差多少 = 姿态估歪多少", 12.5, False, GRAY)],
    [("▸ 小角度时 sinθe ≈ θe → e 就是现成的", 12.5, False, GRAY)],
    [("   线性误差信号,直接当角速度用", 12.5, True, BLUE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
# 右上: v 从哪来
box(s, 6.5, 1.35, 6.4, 2.0, [
    [("v̂ 的来路:把世界系 z 轴转到机体系", 13.5, True, NAVY)],
    [("v̂ = R(q)ᵀ · [0,0,1] ,展开就是代码里的:", 12.5, False, GRAY)],
    [("vx = 2(q1q3−q0q2)  vy = 2(q0q1+q2q3)", 12.5, True, BLUE)],
    [("vz = q0²−q1²−q2²+q3²   (代码存一半值)", 12.5, True, BLUE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
# 右下: PI 物理意义
box(s, 6.5, 3.55, 6.4, 3.35, [
    [("PI 两项各自在干什么", 14, True, NAVY)],
    [("", 3, False, GRAY)],
    [("Kp·e:比例拽回", 13.5, True, BLUE)],
    [("▸ 等价于一阶互补滤波的截止频率:", 12.5, False, GRAY)],
    [("   Kp 大 = 更信加计、收敛快;Kp 小 = 更信陀螺", 12.5, False, GRAY)],
    [("", 5, False, GRAY)],
    [("Ki·∫e:其实是在线估计陀螺零偏!", 13.5, True, BLUE)],
    [("▸ 恒定的误差被积分器慢慢攒下来,", 12.5, False, GRAY)],
    [("   正好就是'没校干净的零偏',自动补偿", 12.5, False, GRAY)],
    [("", 5, False, GRAY)],
    [("局限:重力绕 z 轴对称 → 加计对 yaw 无感", 12.5, True, ORANGE)],
    [("→ yaw 只靠陀螺,会慢漂(想不漂要磁力计)", 12.5, False, ORANGE)],
], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "回答'叉积凭什么能当误差':几何上它同时给了'歪了多少'和'往哪边转回去'。右下是升华:Ki 项本质是在线零偏估计器——这就是为什么 Mahony 对校准残差不敏感。yaw 不可观也要交代。")

# ============ P13 数值实现细节 ============
s = add_slide()
title_bar(s, "深挖② 数值实现细节", "从公式到代码的四件事")
box(s, 0.55, 1.3, 6.15, 2.9, [
    [("① 离散积分的四行标量式", 14, True, NAVY)],
    [("", 3, False, GRAY)],
    [("先算 g = ω'·dt/2 ,然后(用同一拍旧 q!):", 12.5, False, GRAY)],
    [("q0 += −q1·gx − q2·gy − q3·gz", 13, True, BLUE)],
    [("q1 +=  q0·gx + q3·gy − q2·gz", 13, True, BLUE)],
    [("q2 +=  q3·gx − q0·gy + q1·gz", 13, True, BLUE)],
    [("q3 += −q2·gx + q1·gy + q0·gz", 13, True, BLUE)],
    [("▸ 四行必须用更新前的旧 q (先备份)", 12, True, RED)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 0.55, 4.4, 6.15, 2.5, [
    [("② 每步归一化,为什么", 14, True, NAVY)],
    [("▸ 截断误差让 |q| 慢慢 ≠ 1", 12.5, False, GRAY)],
    [("▸ |q|≠1 的四元数不再是纯旋转,", 12.5, False, GRAY)],
    [("   姿态会'缩水/变形'", 12.5, False, GRAY)],
    [("▸ q ← q / |q| ,一次 sqrt + 4 次除法", 12.5, True, BLUE)],
    [("   500Hz 下开销可忽略", 12.5, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 6.95, 1.3, 5.95, 2.9, [
    [("③ 欧拉角从旋转矩阵元素来", 14, True, NAVY)],
    [("▸ atan2 不是拍脑袋:正好取自", 12.5, False, GRAY)],
    [("   R(q) 里对应的两两组合", 12.5, False, GRAY)],
    [("▸ 单值域:roll/yaw ∈ (−π,π]", 12.5, False, GRAY)],
    [("   pitch ∈ [−π/2, π/2] 用 asin", 12.5, False, GRAY)],
    [("▸ asin 内层必须钳位,否则 NaN", 12.5, True, RED)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 6.95, 4.4, 5.95, 2.5, [
    [("④ 三个高频翻车点(实测踩过)", 14, True, RED)],
    [("▸ 单位混乱:deg/s 与 rad/s 混用 →", 12.5, False, GRAY)],
    [("   姿态慢 57 倍,一眼假", 12.5, False, GRAY)],
    [("▸ 忘减零偏 → 静止也在'匀速旋转'", 12.5, False, GRAY)],
    [("▸ dt 用猜不用测 → 收敛速度不对", 12.5, False, GRAY)],
], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "公式到代码的最后一公里:四行标量式与仓库 mahony.c 逐行对应;'旧q'是经典错;归一化和asin钳位都是防 NaN 的防线。三个翻车点全是我们真实踩过的。")

# ============ P14 垂直加速度 az ============
s = add_slide()
title_bar(s, "垂直加速度 az", "去掉重力,才知道悬停了没")
box(s, 0.55, 1.3, 6.1, 2.9, [
    [("a_world = R(q) · a_body", 18, True, BLUE)],
    [("", 6, False, GRAY)],
    [("取世界系 z 分量(R 第三行):", 14, False, GRAY)],
    [("az = 2(q1q3−q0q2)·ax + 2(q2q3+q0q1)·ay", 14, True, GRAY)],
    [("     + (q0²−q1²−q2²+q3²)·az_b − g", 14, True, GRAY)],
    [("", 6, False, GRAY)],
    [("物理含义:扣除重力后的净垂直加速度", 14, True, NAVY)],
], fill=LIGHT, line=BLUE)
text_in(s, 0.55, 4.5, 6.1, 2.4, [
    [("为什么需要它?", 16, True, NAVY)],
    [("▸ 没有气压计/高度计 → 无法直接定高", 14, False, GRAY)],
    [("▸ 但'悬停'的物理条件是 az = 0(升力=重力)", 14, False, GRAY)],
    [("▸ 用加速度闭环逼近悬停,不积分不发散", 14, False, GRAY)],
])
box(s, 7.1, 1.3, 5.8, 2.6, [
    [("精度依赖姿态精度", 15, True, NAVY)],
    [("▸ 姿态误差 θ → az 误差 ≈ g·θ²/2(二阶小量)", 13.5, False, GRAY)],
    [("▸ 小角度下安全;大倾角需姿态准", 13.5, False, GRAY)],
    [("▸ 曾发现系数符号写反(转置矩阵)", 13.5, False, ORANGE)],
    [("   → 已修复并加回归用例", 13.5, False, GREEN)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 7.1, 4.2, 5.8, 2.8, [
    [("az 环怎么用", 15, True, NAVY)],
    [("悬停态:throttle −= az × 0.0002 (每 5ms)", 14, True, BLUE)],
    [("az>0 在升 → 收油门;az<0 在掉 → 加油门", 13.5, False, GRAY)],
    [("收敛到 az≈0,即'加速度意义下的悬停'", 13.5, False, GRAY)],
    [("", 6, False, GRAY)],
    [("局限:速度不为零时托住的是'趋势',会缓慢漂", 12.5, False, ORANGE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "需求反推传感器:想要定高→没有高度计→退而求其次用加速度。坦诚局限。转置矩阵坑在此重提,呼应'对着标准式逐项核对'。")

# ============ P15 Mahony 实用要点 ============
s = add_slide()
title_bar(s, "姿态解算实用要点", "门限 · 调参 · 滤波")
box(s, 0.55, 1.35, 4.0, 5.5, [
    [("修正门限(重要)", 15, True, RED)],
    [("", 4, False, GRAY)],
    [("当 |a| ∉ [0.5, 20] m/s²:", 13.5, True, RED)],
    [("跳过加计修正,只信陀螺", 13.5, True, RED)],
    [("", 6, False, GRAY)],
    [("为什么?", 14, True, NAVY)],
    [("▸ 机动/振动时,加计测的不只是重力", 13, False, GRAY)],
    [("▸ 这时它的'倾角信息'是错的", 13, False, GRAY)],
    [("▸ 宁可暂时漂,也不能被它带偏", 13, False, GRAY)],
    [("", 8, False, GRAY)],
    [("→ 类比:导航员晕船时,先别听他的", 13, True, BLUE)],
], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 4.75, 1.35, 4.0, 5.5, [
    [("两个参数怎么调", 15, True, NAVY)],
    [("", 4, False, GRAY)],
    [("Kp = 0.5:拽回的力度", 14.5, True, BLUE)],
    [("▸ 大:收敛快但姿态抖", 13, False, GRAY)],
    [("▸ 小:平稳但纠偏慢", 13, False, GRAY)],
    [("▸ 判据:静止放置 1~2 秒对准水平", 13, True, GRAY)],
    [("", 8, False, GRAY)],
    [("Ki = 0.05:消'慢性子'偏差", 14.5, True, BLUE)],
    [("▸ 针对陀螺残余零偏", 13, False, GRAY)],
    [("▸ 从 Kp/10 起步", 13, False, GRAY)],
    [("", 8, False, GRAY)],
    [("本项目: Kp=0.5, Ki=0.05", 13, True, GREEN)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 8.95, 1.35, 3.95, 5.5, [
    [("进数据前先滤波", 15, True, NAVY)],
    [("", 4, False, GRAY)],
    [("y[n]=(1−α)y[n−1]+αx[n]", 13.5, True, BLUE)],
    [("", 6, False, GRAY)],
    [("▸ 电机/机架振动几十~几百Hz", 13, False, GRAY)],
    [("▸ 姿态信号 <20Hz,滤掉不吃亏", 13, False, GRAY)],
    [("▸ 两级级联=二阶,更干净", 13, False, GRAY)],
    [("", 8, False, GRAY)],
    [("取舍", 14, True, NAVY)],
    [("▸ 滤得越狠,相位滞后越大", 13, False, ORANGE)],
    [("▸ 加计重滤(只供慢基准)", 13, False, GRAY)],
    [("▸ 陀螺轻滤(内环要快)", 13, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "三个实操问题:门限、调参判据、滤波取舍——看完原理后真正上手会问的三件事。")

# ============ P16 系统架构 ============
s = add_slide()
title_bar(s, "系统架构", "四任务流水线 + 中断前端")
box(s, 0.5, 1.45, 1.75, 0.75, [("MPU6050", 14, True, GRAY), ("I²C 400kHz", 11, False, GRAY)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
box(s, 0.5, 2.9, 1.75, 0.9, [("手机 App", 14, True, GRAY), ("BLE→UART", 11, False, GRAY)], fill=RGBColor(0xDE,0xEB,0xF7), line=BLUE)
box(s, 3.0, 1.3, 2.0, 1.05, [("ImuRead 任务", 14, True, NAVY), ("500Hz · 优先级6", 11.5, False, GRAY), ("采集+零偏校准", 11.5, False, GRAY)], fill=LIGHT, line=BLUE)
box(s, 6.0, 1.3, 2.0, 1.05, [("Attitude 任务", 14, True, NAVY), ("500Hz · 优先级5", 11.5, False, GRAY), ("滤波+Mahony", 11.5, False, GRAY)], fill=LIGHT, line=BLUE)
box(s, 9.0, 1.3, 2.0, 1.05, [("Control 任务", 14, True, NAVY), ("200Hz · 优先级4", 11.5, False, GRAY), ("安全链+串级PID", 11.5, False, GRAY)], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE)
box(s, 11.35, 1.28, 1.55, 1.1, [("Mixer", 13, True, NAVY), ("+ TIM3 PWM", 11, False, GRAY), ("4kHz", 11, False, GRAY)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
box(s, 6.0, 3.1, 2.0, 1.0, [("Comm 任务", 14, True, NAVY), ("事件驱动 · 优先级3", 11.5, False, GRAY), ("BLE帧解析分发", 11.5, False, GRAY)], fill=RGBColor(0xDE,0xEB,0xF7), line=BLUE)
box(s, 3.0, 3.1, 2.0, 1.0, [("USART1 ISR", 13, True, GRAY), ("逐字节接收", 11.5, False, GRAY)], fill=RGBColor(0xDE,0xEB,0xF7), line=BLUE)
arrow(s, 2.25, 1.82, 3.0, 1.82, GREEN)
arrow(s, 5.0, 1.82, 6.0, 1.82)
arrow(s, 8.0, 1.82, 9.0, 1.82, ORANGE)
arrow(s, 11.0, 1.82, 11.35, 1.82, GREEN)
arrow(s, 2.9, 3.6, 3.0, 3.6, BLUE, dash="dash")
arrow(s, 5.0, 3.6, 6.0, 3.6, BLUE)
arrow(s, 7.0, 3.1, 7.0, 2.35, ORANGE)
text_in(s, 5.15, 1.9, 0.9, 0.35, [("队列×10", 11, True, BLUE)], align=PP_ALIGN.CENTER)
text_in(s, 8.15, 1.9, 0.9, 0.35, [("g_att", 11, True, ORANGE)], align=PP_ALIGN.CENTER)
text_in(s, 5.1, 3.62, 0.9, 0.35, [("队列×64", 11, True, BLUE)], align=PP_ALIGN.CENTER)
text_in(s, 7.12, 2.5, 2.4, 0.6, [("SetTarget/Arm/Emergency", 11, True, ORANGE)])
box(s, 0.5, 4.6, 6.1, 2.3, [
    [("设计要点", 14, True, NAVY)],
    [("▸ 采集/解算/控制/通信四任务解耦,数据流单向", 13, False, GRAY)],
    [("▸ 队列做背压与缓冲,ISR 只做入队(短平快)", 13, False, GRAY)],
    [("▸ 零偏校准放 ImuRead:运行期单一 I²C 拥有者", 13, False, GRAY)],
    [("▸ 姿态共享区:高优先级写+临界区读,免锁", 13, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 6.9, 4.6, 6.0, 2.3, [
    [("频率与负载", 14, True, NAVY)],
    [("▸ 500Hz 采样匹配 MPU6050 输出率", 13, False, GRAY)],
    [("▸ 200Hz 控制 = 5ms,角度环带宽足够", 13, False, GRAY)],
    [("▸ Control 每拍喂 IWDG,卡死 1s 内复位", 13, False, GRAY)],
    [("▸ 浮点全走 FPU,CPU 占用估算 <30%", 13, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "原理讲完落回工程:四任务流水线对应前面四层。重点讲队列解耦和I2C单一拥有者。")

# ============ P17 PID 原理 ============
s = add_slide()
title_bar(s, "PID 控制原理", "P/I/D 三项各自的本质")
box(s, 0.55, 1.3, 6.0, 1.15, [
    [("连续形式", 13.5, True, NAVY)],
    [("u(t) = Kp·e(t) + Ki·∫e(t)dt + Kd·de/dt", 16.5, True, BLUE)],
    [("e = 目标值 − 反馈值", 12.5, False, GRAY)],
], fill=LIGHT, line=BLUE)
box(s, 6.85, 1.3, 6.05, 1.15, [
    [("离散实现(每拍执行)", 13.5, True, NAVY)],
    [("integral += e·dt;  deriv = (e−e_prev)/dt", 14, True, GREEN)],
    [("u = Kp·e + Ki·integral + Kd·deriv", 14, True, GREEN)],
], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
table(s, 0.55, 2.75, 12.35, [
    ["项", "物理直觉", "带来的好处", "副作用(没有免费的午餐)"],
    ["P 比例", "'现在'的误差,立即纠正", "响应快,主干增益", "留稳态误差;过大致震荡"],
    ["I 积分", "'过去'误差的累积,不归零不罢休", "消除稳态误差", "响应慢;积分饱和→超调"],
    ["D 微分", "'未来'趋势的预判,提前刹车", "阻尼,抑制震荡", "放大噪声;打杆瞬间 kick"],
], [1.1, 3.4, 3.3, 4.55], size=13, row_h=0.62)
box(s, 0.55, 5.55, 6.0, 1.5, [
    [("工程整定顺序", 14.5, True, NAVY)],
    [("① 先 P:加到临界震荡,回退 ~30%", 13, False, GRAY)],
    [("② 再 D:压住震荡(加阻尼)", 13, False, GRAY)],
    [("③ 最后小 I:只用来消静差", 13, False, GRAY)],
], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 6.85, 5.55, 6.05, 1.5, [
    [("两个进阶细节(实现中落地)", 14.5, True, NAVY)],
    [("▸ D 用测量值微分 → 消 derivative kick", 13, False, GRAY)],
    [("▸ I 用输出限幅 + 回算 → 抗积分饱和", 13, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "'现在/过去/未来'讲本质;副作用列是工程取舍;右下两个细节衔接下页实现。")

# ============ PID 闭环框图 ============
s = add_slide()
title_bar(s, "PID 控制逻辑(图解)", "一个误差,三条支路,合成一个输出")
box(s, 0.55, 2.4, 1.5, 0.8, [("目标角", 13.5, True, GRAY), ("±30°", 11.5, False, GRAY)], fill=RGBColor(0xDE,0xEB,0xF7), line=BLUE)
box(s, 2.65, 2.4, 0.55, 0.8, [("⊗", 16, True, NAVY)], fill=WHITE, line=NAVY)
text_in(s, 2.35, 2.05, 1.2, 0.3, [("e = 目标−反馈", 11, True, RED)], align=PP_ALIGN.CENTER)
box(s, 3.95, 1.3, 2.5, 0.75, [("P 支路", 12.5, True, NAVY), ("Kp · e  ('现在')", 12, True, BLUE)], fill=LIGHT, line=BLUE)
box(s, 3.95, 2.42, 2.5, 0.75, [("I 支路", 12.5, True, NAVY), ("Ki · ∫e dt  ('过去')", 12, True, GREEN)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
box(s, 3.95, 3.55, 2.5, 0.75, [("D 支路", 12.5, True, NAVY), ("Kd · de/dt  ('未来')", 12, True, ORANGE)], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE)
box(s, 7.15, 2.4, 0.55, 0.8, [("Σ", 16, True, NAVY)], fill=WHITE, line=NAVY)
box(s, 8.3, 2.4, 1.9, 0.8, [("执行器", 12.5, True, NAVY), ("混控+电机", 11.5, False, GRAY)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
box(s, 10.75, 2.4, 2.1, 0.8, [("被控对象", 12.5, True, NAVY), ("机身(转动惯量)", 11.5, False, GRAY)], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE)
arrow(s, 2.05, 2.8, 2.65, 2.8, BLUE)
arrow(s, 3.2, 2.8, 3.6, 1.67); gline(s, 3.6, 1.67, 3.95, 1.67, BLUE, 2)
gline(s, 3.2, 2.8, 3.95, 2.8, BLUE, 2)
arrow(s, 3.2, 2.8, 3.6, 3.92); gline(s, 3.6, 3.92, 3.95, 3.92, BLUE, 2)
arrow(s, 6.45, 1.67, 6.9, 1.67); gline(s, 6.9, 1.67, 6.9, 2.8, BLUE, 2)
arrow(s, 6.45, 2.8, 7.15, 2.8, BLUE)
arrow(s, 6.45, 3.92, 6.9, 3.92); gline(s, 6.9, 3.92, 6.9, 2.8, BLUE, 2)
text_in(s, 6.55, 1.25, 1.6, 0.3, [("u = P+I+D", 11, True, NAVY)])
arrow(s, 7.7, 2.8, 8.3, 2.8)
arrow(s, 10.2, 2.8, 10.75, 2.8)
arrow(s, 11.8, 3.2, 11.8, 4.1, NAVY)
text_in(s, 11.15, 4.12, 1.5, 0.3, [("姿态角 θ", 11.5, True, NAVY)])
box(s, 8.3, 4.55, 3.5, 0.75, [("传感器", 12.5, True, NAVY), ("IMU + Mahony(前两章)", 11.5, False, GRAY)], fill=LIGHT, line=BLUE)
gline(s, 11.8, 4.1, 11.8, 4.9, NAVY, 2); arrow(s, 11.8, 4.9, 11.85, 4.9, NAVY, 2)
gline(s, 11.8, 4.9, 2.9, 4.9, NAVY, 2)
arrow(s, 2.9, 4.9, 2.9, 3.2, NAVY)
text_in(s, 3.1, 4.95, 5.5, 0.3, [("负反馈:反馈量与目标相减", 11.5, True, NAVY)])
box(s, 0.55, 5.6, 6.0, 1.35, [
    [("负反馈在干什么", 14, True, NAVY)],
    [("▸ e>0(没到目标)→ u>0 → 姿态朝 e 减小的方向动", 12.5, False, GRAY)],
    [("▸ e<0(冲过头)→ u<0 → 往回拉 → 自动收敛", 12.5, False, GRAY)],
], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 6.85, 5.6, 6.0, 1.35, [
    [("映射到这台无人机", 14, True, NAVY)],
    [("▸ 执行器=混控+电机(受力章);对象=机身惯量", 12.5, False, GRAY)],
    [("▸ 传感器=IMU+Mahony;串级=这个框图套两层", 12.5, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "框图从左到右走一遍:目标与反馈相减得误差,三条支路各干各的,合成 u 去推动电机。强调负反馈的'自动收敛'本质,以及每个方块在本机上的实物对应。")

# ============ 参数对阶跃响应的影响 ============
s = add_slide()
title_bar(s, "Kp / Ki / Kd 对阶跃响应的影响", "看曲线调参(蓝=小,绿=适中,红=大)")

def plot_axes(px, py, w, h, title):
    box(s, px, py, w, h, [], fill=WHITE, line=BORD)
    text_in(s, px + 0.08, py + 0.05, w - 0.2, 0.3, [(title, 12.5, True, NAVY)])
    ax0, ay0 = px + 0.45, py + h - 0.35      # 原点
    ax1, ay1 = px + w - 0.15, py + 0.4       # 右上
    gline(s, ax0, ay1, ax0, ay0, GRAY, 1.5)
    arrow(s, ax0, ay0, ax1, ay0, GRAY, 1.5)
    uy = ay0 - (ay0 - ay1) * 0.62            # y=1 的位置(留出超调空间)
    gline(s, ax0, uy, ax1, uy, GRAY, 1.2, dash="dash")
    text_in(s, ax1 - 0.55, uy - 0.26, 0.6, 0.22, [("目标", 9.5, False, GRAY)])
    return ax0, ax1, ay0, uy

def plot_resp(ax0, ax1, ay0, uy, fn, color, n=40, clip=1.65):
    span = (uy - ay0) / 1.0                  # 1.0 对应的像素高
    pts = []
    for i in range(n + 1):
        t = i / n
        v = fn(t)
        v = max(0.0, min(v, clip))
        pts.append((ax0 + (ax1 - ax0) * t, ay0 - v * span))
    for i in range(len(pts) - 1):
        gline(s, pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1], color, 2.2)

def under_damp(sigma, wd):
    return lambda t: 1 - math.exp(-sigma * t) * (math.cos(wd * t) + (sigma / wd) * math.sin(wd * t))

GW, GH, GY = 3.95, 3.0, 1.35
# Kp 组
a = plot_axes(0.55, GY, GW, GH, "Kp 增大:更快,但震荡↗")
plot_resp(a[0], a[1], a[2], a[3], lambda t: 0.72 * (1 - math.exp(-4.5 * t)), BLUE)          # 小Kp:慢+静差
plot_resp(a[0], a[1], a[2], a[3], under_damp(3.4, 3.6), GREEN)                              # 中:稍过冲
plot_resp(a[0], a[1], a[2], a[3], under_damp(1.1, 9.0), RED)                                # 大:快但震荡
text_in(s, 0.7, GY + GH + 0.08, GW - 0.2, 0.9, [
    [("小:爬得慢,还留着稳态误差", 11.5, True, BLUE)],
    [("大:响应快、误差小,但超调震荡;过大致失稳", 11.5, True, RED)],
])
# Ki 组
b = plot_axes(4.75, GY, GW, GH, "Ki 增大:消静差,但超调↗")
plot_resp(b[0], b[1], b[2], b[3], lambda t: 0.9 * (1 - math.exp(-5 * t)), BLUE)             # 无I:静差0.9
plot_resp(b[0], b[1], b[2], b[3], under_damp(2.6, 4.2), GREEN)                              # 适中
plot_resp(b[0], b[1], b[2], b[3], under_damp(0.8, 7.5), RED)                                # 过大:来回摆
text_in(s, 4.9, GY + GH + 0.08, GW - 0.2, 0.9, [
    [("无:永远差一点(重心偏→悬停歪着)", 11.5, True, BLUE)],
    [("大:消差猛但刹车不住,来回摆", 11.5, True, RED)],
])
# Kd 组
c = plot_axes(8.95, GY, GW, GH, "Kd 增大:阻尼,超调↓")
plot_resp(c[0], c[1], c[2], c[3], under_damp(1.4, 8.0), BLUE)                               # 小D:超调大
plot_resp(c[0], c[1], c[2], c[3], under_damp(3.8, 3.2), GREEN)                              # 适中
plot_resp(c[0], c[1], c[2], c[3], lambda t: (1 - math.exp(-2.2 * t)) + 0.035 * math.exp(-2.5 * t) * math.sin(55 * t), RED)  # 过大:迟钝+噪声毛刺
text_in(s, 9.1, GY + GH + 0.08, GW - 0.2, 0.9, [
    [("小:冲过头才回,震荡几拍", 11.5, True, BLUE)],
    [("大:过于敏感,噪声被微分放大(红色毛刺)", 11.5, True, RED)],
])
box(s, 0.55, 5.75, 12.35, 1.2, [
    [("读图整定:先加 Kp 到'快而略震'(绿) → 加 Kd 把超调压掉 → 最后小 Ki 收静差", 14, True, NAVY)],
    [("口诀:P 定快慢,D 定稳不稳,I 定准不准;一项一项动,每次只动一个", 13, False, GRAY)],
], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "三张图是自控课的看家图:同一条二阶欠阻尼曲线,改变阻尼/增益/积分就得到不同形态。讲法:先指目标虚线,再讲蓝色缺什么、红色多什么、绿色为什么合适。底部读图整定与上一页口诀闭环。")

# ============ P18 串级实现 ============
s = add_slide()
title_bar(s, "串级 PID 实现", "角度环 P + 角速度环 PID")
box(s, 0.55, 1.35, 1.9, 0.7, [("目标角", 13, True, GRAY), ("±30°", 11.5, False, GRAY)], fill=RGBColor(0xDE,0xEB,0xF7), line=BLUE)
box(s, 3.15, 1.35, 1.9, 0.7, [("角度环 P", 13, True, NAVY), ("×6", 11.5, False, GRAY)], fill=LIGHT, line=BLUE)
box(s, 5.75, 1.35, 1.9, 0.7, [("期望角速度", 13, True, GRAY), ("限±200°/s", 11.5, False, GRAY)], fill=RGBColor(0xDE,0xEB,0xF7), line=BLUE)
box(s, 8.35, 1.35, 1.9, 0.7, [("角速度环 PID", 13, True, NAVY), ("200Hz", 11.5, False, GRAY)], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE)
box(s, 10.95, 1.35, 1.9, 0.7, [("力矩 ±", 13, True, GRAY), ("R/P/Y", 11.5, False, GRAY)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
for xx in (2.45, 5.05, 7.65, 10.25):
    arrow(s, xx, 1.7, xx + 0.7, 1.7)
text_in(s, 3.0, 2.15, 2.2, 0.4, [("e = 目标角−姿态角", 11, False, BLUE)], align=PP_ALIGN.CENTER)
text_in(s, 8.1, 2.15, 2.4, 0.4, [("e = 期望−实测ω", 11, False, ORANGE)], align=PP_ALIGN.CENTER)
table(s, 0.55, 2.75, 6.0, [
    ["参数", "值", "作用"],
    ["kp_angle", "6.0", "角度差→期望角速度(姿态自稳刚度)"],
    ["kp_rate", "0.15", "角速度误差→力矩(阻尼)"],
    ["ki_rate", "2.0 /s", "消除恒值扰动(重心偏/风)"],
    ["kd_rate", "0.005", "角速度变化率的阻尼(抑震荡)"],
], [1.3, 0.9, 3.8], size=12.5, row_h=0.46)
box(s, 0.55, 5.25, 6.0, 1.7, [
    [("为什么串级?", 14.5, True, NAVY)],
    [("▸ 单级角度PID:姿态已经错了才大力纠", 13, False, GRAY)],
    [("▸ 内环以200Hz锁死角速度,扰动在'变成", 13, False, GRAY)],
    [("   姿态错误'前就被吸收", 13, False, GRAY)],
    [("▸ yaw 无角度外环(无磁力计,航向不可观)", 12.5, False, ORANGE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 6.85, 2.75, 6.0, 4.2, [
    [("三个工程细节(踩坑修复项)", 14.5, True, RED)],
    [("", 4, False, GRAY)],
    [("① 积分 ×dt:离散积分量纲归一", 13.5, True, GRAY)],
    [("   I_out = ki·Σ(e·dt),不随控制频率变义", 12.5, False, GRAY)],
    [("", 4, False, GRAY)],
    [("② 抗饱和 = 输出限幅 + 回算", 13.5, True, GRAY)],
    [("   I 项权限钳 ±0.15(15%油门),超限回写", 12.5, False, GRAY)],
    [("", 4, False, GRAY)],
    [("③ D 项用测量值微分", 13.5, True, GRAY)],
    [("   d(−gyro)/dt 而非 d(err)/dt", 12.5, False, GRAY)],
    [("   → 打杆瞬间无 derivative kick", 12.5, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "呼应动力学页:角速度环的输出就是那三轴力矩。三个工程细节都是真实修过的坑。")

# ============ P19 混控 ============
s = add_slide()
title_bar(s, "X 型混控", "把第二章的力矩公式变成代码")
box(s, 0.7, 1.5, 5.6, 4.4, [
    [("M2 前左(逆)      M1 前右(顺)", 15, True, GRAY)],
    [("      ╲            ╱", 16, False, GRAY)],
    [("        ╲        ╱", 16, False, GRAY)],
    [("           ╳   机头↑", 15, True, NAVY)],
    [("        ╱        ╲", 16, False, GRAY)],
    [("      ╱            ╲", 16, False, GRAY)],
    [("M3 后左(顺)      M4 后右(逆)", 15, True, GRAY)],
], fill=WHITE, line=GRAY)
box(s, 6.9, 1.45, 6.0, 2.5, [
    [("混控矩阵 (T=油门, R/P/Y=力矩)", 14.5, True, NAVY)],
    [("M1 = T + P − R − Y      M2 = T + P + R + Y", 14.5, True, BLUE)],
    [("M3 = T − P + R − Y      M4 = T − P − R + Y", 14.5, True, BLUE)],
    [("各路 clamp [0,1] 后写 TIM3 CCR", 13, False, GRAY)],
], fill=LIGHT, line=BLUE)
text_in(s, 6.9, 4.15, 6.0, 0.5, [("对照受力页:+R=左桨加速(右滚) +P=前桨加速(抬头)", 12.5, False, GREEN)])
box(s, 6.9, 4.75, 6.0, 2.2, [
    [("共模/差模视角(关键)", 14.5, True, NAVY)],
    [("▸ 油门 = 共模:四桨同增减 → 只变总推力", 13, False, GRAY)],
    [("▸ 力矩 = 差模:对角互补,和为零 → 只变力矩", 13, False, GRAY)],
    [("▸ 二者近似正交 → 高度环与姿态环解耦", 13, True, BLUE)],
], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
text_in(s, 0.7, 6.05, 5.6, 0.9, [
    [("⚠ 首飞前必台架验证:通道映射与桨向(拆桨!)", 12.5, True, RED)],
])
notes(s, "这页矩阵就是受力分析页四个公式的直接翻译:差动出力矩、共模出推力。共模/差模正交性呼应三环协同页。")

# ============ P20 状态机 ============
s = add_slide()
title_bar(s, "飞行状态机", "五个状态,核心是'油门归谁管'")
sts = [
    ("IDLE", "上电/急停后", "电机停", RGBColor(0xED,0xED,0xED)),
    ("DISARMED", "已解锁,油门<5%", "电机停,等推杆", RGBColor(0xDE,0xEB,0xF7)),
    ("FOLLOW", "油门>5%", "油门直接跟杆", RGBColor(0xE2,0xEF,0xDA)),
    ("HOVER", "松杆", "az闭环悬停", RGBColor(0xFD,0xF0,0xE6)),
    ("LANDING", "收油/超时", "匀速缓降-0.3/s", RGBColor(0xDE,0xEB,0xF7)),
]
x = 0.55
for name, cond, act, fillc in sts:
    box(s, x, 1.5, 2.25, 1.35, [
        [(name, 16, True, NAVY)], [(cond, 11.5, False, GRAY)], [(act, 11.5, False, GRAY)]
    ], fill=fillc, line=BLUE)
    if x < 10.6:
        arrow(s, x + 2.25, 2.17, x + 2.5, 2.17)
    x += 2.5
arrow(s, 12.9, 2.9, 1.7, 2.9, color=ORANGE, dash="dash")
arrow(s, 1.7, 2.9, 1.7, 2.85, color=ORANGE, dash="dash")
text_in(s, 4.5, 2.95, 5.0, 0.4, [("落地自动上锁 → IDLE(任意时刻可被急停拉回)", 12, True, ORANGE)], align=PP_ALIGN.CENTER)
box(s, 0.55, 3.7, 6.0, 3.2, [
    [("迁移规则", 14.5, True, NAVY)],
    [("▸ 解锁+推油门>5% → FOLLOW 起飞", 13.5, False, GRAY)],
    [("▸ 松杆 → HOVER:保留当前油门为闭环起点", 13.5, False, GRAY)],
    [("▸ HOVER 满 3s → 自动缓降(独立计时器)", 13.5, False, GRAY)],
    [("▸ 再推杆 → 从 HOVER 抢回 FOLLOW", 13.5, False, GRAY)],
    [("▸ 落地(油门降到0) → 自动上锁回 IDLE", 13.5, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 6.85, 3.7, 6.0, 3.2, [
    [("操作者视角", 14.5, True, NAVY)],
    [("▸ 姿态杆任何飞行状态都有效(姿态环常开)", 13.5, False, GRAY)],
    [("▸ 唯一时序规则:'松杆3秒必降落'", 13.5, True, BLUE)],
    [("▸ 悬停中打杆平移,az环自动补倾角掉高", 13.5, False, GRAY)],
    [("▸ 落地自动 disarm,不会'落地又窜起'", 13.5, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "状态机一句话:油门归属权切换,姿态永远归PID。")

# ============ P21 安全链 ============
s = add_slide()
title_bar(s, "七级安全链", "每 5ms 自上而下短路评估")
table(s, 0.55, 1.3, 12.3, [
    ["优先级", "检查项", "触发条件", "动作", "设计考量"],
    ["1", "急停锁存", "emergency 标志", "电机立即停", "上锁才能清除,防误触恢复"],
    ["2", "未解锁", "!armed", "电机停", "上电锁:双重防线"],
    ["2.5", "数据有效性", "NaN/Inf", "立即急停", "NaN 比较恒假,会绕过一切保护"],
    ["3", "姿态超限", "|角|>45° 持续200ms", "急停", "去抖防单帧噪声误杀"],
    ["4", "失控", "1s 无指令", "转缓降", "不直接停机(半空断电=摔)"],
    ["5", "悬停超时", "HOVER 满 3s", "转缓降", "无高度计,限制久悬"],
    ["6", "姿态丢失", "0.5s 无新姿态", "急停", "防电机冻结在最后占空比"],
    ["7", "看门狗", "IWDG 1s 未喂", "芯片复位", "任务卡死兜底"],
], [0.9, 1.7, 2.5, 1.5, 5.7], size=12.5, row_h=0.5)
box(s, 0.55, 5.85, 12.3, 1.05, [
    [("设计原则:失效安全逐级降级 —— 能缓降不悬停,能悬停不停机,该停机绝不犹豫", 14.5, True, NAVY)],
], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE)
notes(s, "每条保护都有'为什么':NaN绕过一切if、单帧噪声不能停机、半空断电比缓降危险。")

# ============ P22 三环协同 ============
s = add_slide()
title_bar(s, "三环如何协同", "呼应受力分析:推力方向 vs 大小")
box(s, 0.55, 1.35, 5.9, 1.05, [("az 环(慢环):管推力大小", 14, True, NAVY), ("积分式调油门 ≤0.4油门/s", 12, False, GRAY)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
box(s, 0.55, 3.35, 5.9, 1.05, [("姿态环→角速度环(快环):管推力方向", 13.5, True, NAVY), ("200Hz,毫秒级", 12, False, GRAY)], fill=RGBColor(0xFD,0xF0,0xE6), line=ORANGE)
box(s, 7.0, 2.3, 2.3, 1.15, [("Mixer", 15, True, NAVY), ("共模‖差模", 12.5, False, GRAY)], fill=LIGHT, line=BLUE)
box(s, 9.9, 2.3, 2.9, 1.15, [("4×电机", 15, True, GRAY), ("PWM 4kHz", 12.5, False, GRAY)], fill=RGBColor(0xE2,0xEF,0xDA), line=GREEN)
arrow(s, 3.5, 2.4, 7.0, 2.6, GREEN); arrow(s, 3.5, 3.9, 7.0, 3.2, ORANGE)
arrow(s, 9.3, 2.87, 9.9, 2.87, GREEN)
text_in(s, 5.7, 1.95, 1.6, 0.4, [("油门T", 11.5, True, GREEN)])
text_in(s, 5.35, 4.0, 1.9, 0.4, [("力矩R/P/Y", 11.5, True, ORANGE)])
box(s, 0.55, 5.0, 12.35, 1.95, [
    [("三条真实耦合通路", 14.5, True, RED)],
    [("① 倾角→升力损失:垂直升力×cosθ → az环被动补油门(悬停压坡不掉高,有利)", 12.5, False, GRAY)],
    [("② 电机饱和抢权限:油门贴 0.85/0.30 边界时姿态差动被削 → 预留15%余量(待做去饱和)", 12.5, False, GRAY)],
    [("③ 估计误差串扰:az 由姿态四元数算出,姿态不准→az不准(二阶小量)", 12.5, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "把开场受力分析闭环:姿态环管'推力指哪',az环管'推力多大',混控的正交性让两者近似不打架。")

# ============ P23 通信协议 ============
s = add_slide()
title_bar(s, "通信协议设计", "帧五要素 + 逐字节状态机")
fr = [("0xA5", "STX 定界", NAVY), ("CMD", "命令字", BLUE), ("LEN", "长度", BLUE),
      ("DATA 0~16B", "小端float32", GRAY), ("XOR", "校验", ORANGE), ("0x5A", "ETX", NAVY)]
x = 0.7
for t, d, c in fr:
    box(s, x, 1.45, 1.85, 0.95, [[(t, 14.5, True, WHITE)], [(d, 11, False, WHITE)]], fill=c, line=c)
    x += 2.05
box(s, 0.7, 2.85, 6.0, 3.6, [
    [("帧设计原则(任何二进制协议通用)", 14.5, True, NAVY)],
    [("① 定界符:流中找帧头(允许丢字节重同步)", 12.5, False, GRAY)],
    [("② 长度:先读长度再收数据,防粘包", 12.5, False, GRAY)],
    [("③ 校验:XOR/CRC 检错,失败整帧丢弃", 12.5, False, GRAY)],
    [("④ 帧尾:独立校验,防'校验值恰好=帧尾'", 12.5, False, GRAY)],
    [("⑤ 数据带范围与 isfinite 校验", 12.5, False, GRAY)],
    [("", 8, False, GRAY)],
    [("ISR 只入队,任务里解析(中断要短)", 13, True, BLUE)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
box(s, 7.1, 2.85, 5.8, 3.6, [
    [("逐字节状态机", 14.5, True, NAVY)],
    [("IDLE → CMD → LEN → DATA → CHK → ETX", 13, True, BLUE)],
    [("", 6, False, GRAY)],
    [("▸ 任何一步非法 → 回 IDLE 重新找帧头", 12.5, False, GRAY)],
    [("▸ 天然处理:断流/粘包/半帧/误码", 12.5, False, GRAY)],
    [("", 8, False, GRAY)],
    [("常用命令示例", 13.5, True, NAVY)],
    [("0x40 SETPOINT: 4×float (T/roll/pitch/yaw)", 12.5, False, GRAY)],
    [("0x24 ARM: 1字节   0x20: 急停", 12.5, False, GRAY)],
], fill=LIGHT, line=BLUE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
notes(s, "协议五要素是通用方法论;状态机解析可展开讲重同步。")

# ============ P24 结束页 ============
s = add_slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 4.7, 13.333, 0.05, ORANGE)
text_in(s, 1.0, 2.1, 11.3, 1.0, [("Q & A", 44, True, WHITE)])
text_in(s, 1.0, 3.35, 11.3, 0.7, [("受力与扭矩 · IMU · 姿态解算 · 控制 —— 每一块都可以往深了聊", 18, False, RGBColor(0xBD,0xD7,0xEE))])
text_in(s, 1.0, 5.0, 11.3, 1.2, [
    [("谢谢聆听", 20, True, WHITE)],
    [("代码: github.com/xiexiaoan147-cyber/STM32F4_Drone", 14, False, RGBColor(0x8E,0xAA,0xDB))],
])
notes(s, "收尾:主动开放提问方向——叉积误差的几何、四元数积分、串级整定都可以聊。")

out = "/home/user/STM32F4_Drone/Docs/四轴飞控原理分享.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("saved:", out, os.path.getsize(out), "bytes")
